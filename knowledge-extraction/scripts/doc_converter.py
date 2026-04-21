#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
企业知识萃取 - 文档批量转换工具
支持 DOCX/XLSX/PPTX → MD 批量转换
"""

import os
import sys
import json
import hashlib
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional

try:
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation
except ImportError as e:
    print(f"错误: 缺少必要的库。请运行: pip install python-docx openpyxl python-pptx")
    print(f"详细错误: {e}")
    sys.exit(1)


class DocumentConverter:
    """文档批量转换器"""

    def __init__(self, source_dir: str, output_dir: str = None):
        """
        初始化转换器

        Args:
            source_dir: 源文档目录
            output_dir: 输出目录（默认为 source_dir/converted）
        """
        self.source_dir = Path(source_dir)
        self.output_dir = Path(output_dir) if output_dir else self.source_dir / "converted"
        self.log_file = self.output_dir / "conversion_log.json"

        # 创建输出目录
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # 加载之前的转换日志
        self.conversion_history = self._load_conversion_log()

        # 支持的文件格式
        self.supported_formats = {
            '.docx': self._convert_docx,
            '.xlsx': self._convert_xlsx,
            '.pptx': self._convert_pptx,
        }

    def _load_conversion_log(self) -> Dict:
        """加载转换日志"""
        if self.log_file.exists():
            try:
                with open(self.log_file, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except:
                return {}
        return {}

    def _save_conversion_log(self):
        """保存转换日志"""
        with open(self.log_file, 'w', encoding='utf-8') as f:
            json.dump(self.conversion_history, f, ensure_ascii=False, indent=2)

    def _get_file_hash(self, file_path: Path) -> str:
        """计算文件哈希值"""
        with open(file_path, 'rb') as f:
            return hashlib.md5(f.read()).hexdigest()

    def _is_converted(self, file_path: Path) -> bool:
        """检查文件是否已转换"""
        file_key = str(file_path)
        if file_key not in self.conversion_history:
            return False

        # 检查文件是否修改过
        current_hash = self._get_file_hash(file_path)
        return self.conversion_history[file_key]['hash'] == current_hash

    def _mark_converted(self, file_path: Path, output_file: Path):
        """标记文件已转换"""
        file_key = str(file_path)
        self.conversion_history[file_key] = {
            'hash': self._get_file_hash(file_path),
            'output': str(output_file),
            'timestamp': datetime.now().isoformat()
        }

    def _convert_docx(self, file_path: Path) -> str:
        """转换 Word 文档为 Markdown"""
        try:
            doc = Document(file_path)
            md_content = []

            # 提取段落
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    # 简单处理标题（根据样式名称）
                    style_name = para.style.name if para.style else ""
                    if 'Heading 1' in style_name or '标题 1' in style_name:
                        md_content.append(f"# {text}\n")
                    elif 'Heading 2' in style_name or '标题 2' in style_name:
                        md_content.append(f"## {text}\n")
                    elif 'Heading 3' in style_name or '标题 3' in style_name:
                        md_content.append(f"### {text}\n")
                    else:
                        md_content.append(f"{text}\n")

            # 提取表格
            for table in doc.tables:
                md_content.append("\n**表格**\n")
                for i, row in enumerate(table.rows):
                    row_data = [cell.text.strip() for cell in row.cells]
                    if i == 0:
                        # 表头
                        md_content.append("| " + " | ".join(row_data) + " |")
                        md_content.append("| " + " | ".join(["---"] * len(row_data)) + " |")
                    else:
                        # 表格内容
                        md_content.append("| " + " | ".join(row_data) + " |")
                md_content.append("\n")

            return "\n".join(md_content)
        except Exception as e:
            return f"转换失败: {str(e)}"

    def _convert_xlsx(self, file_path: Path) -> str:
        """转换 Excel 文档为 Markdown"""
        try:
            wb = load_workbook(file_path, read_only=True, data_only=True)
            md_content = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                md_content.append(f"\n## 工作表: {sheet_name}\n")

                # 获取有数据的区域
                for row in ws.iter_rows(values_only=True):
                    # 过滤空行
                    if any(cell is not None for cell in row):
                        row_data = [str(cell) if cell is not None else "" for cell in row]
                        md_content.append("| " + " | ".join(row_data) + " |")

                if ws.max_row > 0:
                    # 添加表头分隔线
                    first_row_length = len(list(ws.iter_rows(min_row=1, max_row=1, values_only=True))[0])
                    md_content.insert(len(md_content) - ws.max_row + 1, "|" + "|".join(["---"] * first_row_length) + "|")

            wb.close()
            return "\n".join(md_content)
        except Exception as e:
            return f"转换失败: {str(e)}"

    def _convert_pptx(self, file_path: Path) -> str:
        """转换 PowerPoint 文档为 Markdown"""
        try:
            prs = Presentation(file_path)
            md_content = []

            for slide_num, slide in enumerate(prs.slides, 1):
                md_content.append(f"\n## 幻灯片 {slide_num}\n")

                # 提取文本框
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        # 根据字体大小判断是否为标题
                        if hasattr(shape, "text_frame"):
                            for paragraph in shape.text_frame.paragraphs:
                                if paragraph.text.strip():
                                    # 简单的标题判断逻辑
                                    font_size = paragraph.font.size if paragraph.font else None
                                    if font_size and font_size > 280000:  # 约18pt以上
                                        md_content.append(f"### {paragraph.text.strip()}\n")
                                    else:
                                        md_content.append(f"{paragraph.text.strip()}\n")
                        else:
                            md_content.append(f"{text}\n")

            return "\n".join(md_content)
        except Exception as e:
            return f"转换失败: {str(e)}"

    def convert_file(self, file_path: Path) -> Optional[Path]:
        """
        转换单个文件

        Args:
            file_path: 源文件路径

        Returns:
            输出文件路径，如果不需要转换则返回 None
        """
        # 检查文件格式
        ext = file_path.suffix.lower()
        if ext not in self.supported_formats:
            return None

        # 检查是否已转换
        if self._is_converted(file_path):
            print(f"  ✓ 跳过（已转换）: {file_path.name}")
            return None

        print(f"  → 转换中: {file_path.name}")

        # 执行转换
        converter = self.supported_formats[ext]
        md_content = converter(file_path)

        # 保存转换结果
        output_file = self.output_dir / f"{file_path.stem}.md"
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(f"# {file_path.name}\n\n")
            f.write(md_content)

        # 标记已转换
        self._mark_converted(file_path, output_file)

        print(f"  ✓ 完成: {output_file.name}")
        return output_file

    def batch_convert(self, recursive: bool = True) -> Dict[str, int]:
        """
        批量转换目录中的文档

        Args:
            recursive: 是否递归处理子目录

        Returns:
            转换统计信息
        """
        print(f"\n开始批量转换...")
        print(f"源目录: {self.source_dir}")
        print(f"输出目录: {self.output_dir}\n")

        # 扫描文件
        if recursive:
            files = [f for f in self.source_dir.rglob("*") if f.is_file()]
        else:
            files = [f for f in self.source_dir.glob("*") if f.is_file()]

        # 过滤支持的格式
        files_to_convert = [f for f in files if f.suffix.lower() in self.supported_formats]

        stats = {
            'total': len(files_to_convert),
            'converted': 0,
            'skipped': 0,
            'failed': 0,
            'by_format': {}
        }

        # 批量转换
        for file_path in files_to_convert:
            ext = file_path.suffix.lower()
            if ext not in stats['by_format']:
                stats['by_format'][ext] = 0

            try:
                result = self.convert_file(file_path)
                if result:
                    stats['converted'] += 1
                    stats['by_format'][ext] += 1
                else:
                    stats['skipped'] += 1
            except Exception as e:
                print(f"  ✗ 转换失败: {file_path.name} - {str(e)}")
                stats['failed'] += 1

        # 保存转换日志
        self._save_conversion_log()

        # 打印统计信息
        print(f"\n转换完成！")
        print(f"总计: {stats['total']} 个文件")
        print(f"已转换: {stats['converted']} 个")
        print(f"跳过: {stats['skipped']} 个")
        print(f"失败: {stats['failed']} 个")
        print(f"\n按格式统计:")
        for ext, count in stats['by_format'].items():
            print(f"  {ext}: {count} 个")

        return stats


def main():
    """命令行入口"""
    import argparse

    parser = argparse.ArgumentParser(description='企业知识萃取 - 文档批量转换工具')
    parser.add_argument('source_dir', help='源文档目录')
    parser.add_argument('-o', '--output', help='输出目录（默认为源目录/converted）')
    parser.add_argument('-r', '--recursive', action='store_true', default=True,
                        help='递归处理子目录（默认启用）')
    parser.add_argument('--no-recursive', action='store_false', dest='recursive',
                        help='不递归处理子目录')

    args = parser.parse_args()

    # 创建转换器
    converter = DocumentConverter(args.source_dir, args.output)

    # 批量转换
    converter.batch_convert(recursive=args.recursive)


if __name__ == '__main__':
    main()
